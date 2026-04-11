import fitz  # PyMuPDF
from docx import Document
import openpyxl
from pptx import Presentation
import os

def parse_txt(file_path):
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()

def parse_pdf(file_path):
    text = ""
    with fitz.open(file_path) as doc:
        for page in doc:
            text += page.get_text()
    return text

def parse_docx(file_path):
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def parse_xlsx(file_path):
    wb = openpyxl.load_workbook(file_path, data_only=True)
    text = ""
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        for row in ws.iter_rows(values_only=True):
            text += " ".join([str(cell) for cell in row if cell is not None]) + " "
    return text

def parse_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
    return text

def extract_text(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == '.txt':
        return parse_txt(file_path)
    elif ext == '.pdf':
        return parse_pdf(file_path)
    elif ext == '.docx':
        return parse_docx(file_path)
    elif ext == '.xlsx':
        return parse_xlsx(file_path)
    elif ext == '.pptx':
        return parse_pptx(file_path)
    else:
        # Try to read as plain text for unknown formats
        try:
            return parse_txt(file_path)
        except:
            return ""
